"""PowerPoint extraction using python-pptx."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from ingest.extractors.base import ExtractedDocument


class PPTXExtractor:
    """Extract slide titles, body text, and speaker notes from presentations."""

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        """Extract text from every slide and return one section per slide."""

        presentation = Presentation(BytesIO(file_bytes))
        sections: list[dict[str, str]] = []
        all_text: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            title = f"Slide {index}"
            body: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if not body:
                        title = text.splitlines()[0] or title
                    body.append(text)
            notes = self._extract_notes(slide)
            if notes:
                body.append(f"Speaker notes: {notes}")
            content = "\n".join(body).strip()
            sections.append({"heading": title, "content": content})
            if content:
                all_text.append(content)

        raw_text = "\n\n".join(all_text)
        props = presentation.core_properties
        return ExtractedDocument(
            title=props.title or filename,
            raw_text=raw_text,
            metadata={
                "author": props.author or None,
                "date_created": props.created.isoformat() if props.created else None,
                "date_modified": props.modified.isoformat() if props.modified else None,
                "page_count": None,
                "slide_count": len(presentation.slides),
                "word_count": len(raw_text.split()),
            },
            sections=sections,
        )

    def _extract_notes(self, slide: object) -> str:
        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide is None:
            return ""
        notes: list[str] = []
        for shape in notes_slide.notes_text_frame.paragraphs:
            text = shape.text.strip()
            if text:
                notes.append(text)
        return "\n".join(notes)
