"""Extractor tests with synthetic office files."""

from __future__ import annotations

from io import BytesIO

from docx import Document
from pptx import Presentation

from ingest.extractors.docx_extractor import DOCXExtractor
from ingest.extractors.pdf_extractor import PDFExtractor
from ingest.extractors.pptx_extractor import PPTXExtractor


def test_docx_extractor_reads_heading_and_text():
    buffer = BytesIO()
    doc = Document()
    doc.add_heading("Radar Study", level=1)
    doc.add_paragraph("FPGA radar processing content.")
    doc.save(buffer)
    extracted = DOCXExtractor().extract(buffer.getvalue(), "radar.docx")
    assert "FPGA radar" in extracted.raw_text
    assert extracted.sections[0]["heading"] == "Radar Study"


def test_pptx_extractor_reads_slide_text():
    buffer = BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "UAV Navigation"
    slide.placeholders[1].text = "GNSS and inertial fusion"
    prs.save(buffer)
    extracted = PPTXExtractor().extract(buffer.getvalue(), "nav.pptx")
    assert "GNSS" in extracted.raw_text
    assert extracted.metadata["slide_count"] == 1


def test_pdf_extractor_handles_invalid_pdf_gracefully():
    extracted = PDFExtractor().extract(b"not a pdf", "broken.pdf")
    assert extracted.title == "broken.pdf"
    assert extracted.raw_text == ""
