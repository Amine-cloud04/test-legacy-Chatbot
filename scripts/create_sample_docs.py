"""Create a larger French demo corpus for local retrieval and LLM testing."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from docx import Document
from pptx import Presentation


SAMPLE_DIR = Path(__file__).resolve().parents[1] / "sample_docs"


def _clear_existing_samples() -> None:
    """Remove old generated demo documents from sample_docs."""

    if not SAMPLE_DIR.exists():
        return
    for path in SAMPLE_DIR.iterdir():
        if path.is_file() and path.suffix.lower() in {".docx", ".pptx", ".pdf"}:
            path.unlink()


def _repeat(text: str, scale: int) -> str:
    """Expand a paragraph with distinct supporting sentences."""

    base = text.strip()
    if scale <= 1:
        return base

    additions = [
        "Le dossier conserve une traçabilité complète entre exigences, essais et décisions d'architecture.",
        "Les résultats sont archivés avec des critères d'acceptation, des limites connues et des références de test.",
        "Les conclusions détaillent les compromis techniques, les points ouverts et les prochaines actions de validation.",
        "La robustesse est évaluée sur des scénarios nominales et dégradés afin de limiter les effets de bord.",
    ]
    parts = [base]
    for index in range(2, scale + 1):
        suffix = additions[(index - 2) % len(additions)]
        parts.append(f"{base} {suffix}")
    return " ".join(parts)


def _add_docx_sections(document: Document, title: str, sections: list[tuple[str, str]], scale: int) -> None:
    """Write a DOCX document with repeated technical sections."""

    document.core_properties.title = title
    document.core_properties.author = "Assistant de connaissances R&D"
    document.add_heading(title, level=1)
    for heading, paragraph in sections:
        document.add_heading(heading, level=2)
        document.add_paragraph(_repeat(paragraph, scale))
        document.add_paragraph(
            _repeat(
                "Le suivi des risques associe toujours le calendrier de validation, la charge de calcul embarquée, les contraintes de sécurité et la qualité des données de test.",
                scale,
            )
        )
    document.add_heading("Synthèse opérationnelle", level=2)
    document.add_paragraph(
        _repeat(
            "La recommandation globale est de poursuivre les essais sur un banc représentatif, de consolider les sources de données et de conserver une traçabilité documentaire stricte.",
            scale,
        )
    )


def _add_pptx_slides(presentation: Presentation, title: str, slides: list[tuple[str, str]], scale: int) -> None:
    """Write a PPTX presentation with repeated content-rich slides."""

    presentation.core_properties.title = title
    presentation.core_properties.author = "Assistant de connaissances R&D"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Corpus de démonstration pour tester le retrieval et la génération locale"

    for heading, body in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = heading
        slide.placeholders[1].text = _repeat(body, scale)


def create_radar_docx(path: Path, scale: int) -> None:
    """Create a long radar and FPGA project report."""

    document = Document()
    _add_docx_sections(
        document,
        "Étude de faisabilité radar FPGA",
        [
            (
                "Contexte",
                "Le projet radar a évalué une chaîne de traitement embarquée pour la détection à courte portée. "
                "L'équipe a mis en œuvre la compression d'impulsions, le rejet du clutter et le scoring des cibles sur un prototype FPGA.",
            ),
            (
                "Architecture technique",
                "L'architecture utilisait un FPGA Xilinx, un traitement du signal en virgule fixe et un processeur de contrôle léger. "
                "La latence est restée inférieure à 18 millisecondes lors des tests en banc, ce qui convient aux contraintes avioniques.",
            ),
            (
                "Validation",
                "Les essais ont couvert les modes basse consommation, les scénarios de forte fréquence de répétition d'impulsions et les limitations de bande passante mémoire. "
                "Les mesures ont montré que la charge calculatoire restait stable quand le pipeline FPGA était correctement dimensionné.",
            ),
            (
                "Leçons apprises",
                "Le principal risque concernait la bande passante mémoire pendant les fortes fréquences de répétition d'impulsions. "
                "Les travaux futurs compareront l'accélération FPGA à un traitement CPU uniquement pour les modes radar basse consommation.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_radar_validation_docx(path: Path, scale: int) -> None:
    """Create a second radar document focused on test and verification."""

    document = Document()
    _add_docx_sections(
        document,
        "Campagne de validation radar embarqué",
        [
            (
                "Objectifs de test",
                "La campagne a vérifié la stabilité du radar embarqué sur des séquences longues, des changements de température et des profils de cible variables. "
                "Le protocole de test a inclus des seuils de détection, des faux positifs et la répétabilité des mesures.",
            ),
            (
                "Résultats",
                "Les résultats ont confirmé que le filtrage des échos et le rejet du clutter amélioraient la robustesse dans les scènes chargées. "
                "Les performances se dégradaient surtout quand les données d'entrée manquaient de qualité ou quand la mémoire partagée saturait.",
            ),
            (
                "Risques",
                "Les risques principaux sont la saturation des bus mémoire, la dérive des paramètres de seuil et le coût de qualification pour une mise en production avionique.",
            ),
            (
                "Suite",
                "La suite du projet doit combiner instrumentation plus fine, journalisation automatique et comparaison avec une chaîne CPU de référence.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_sensor_docx(path: Path, scale: int) -> None:
    """Create a long inertial and GNSS sensor fusion note."""

    document = Document()
    _add_docx_sections(
        document,
        "Notes de fusion capteurs inertiels et GNSS",
        [
            (
                "Matrice de capteurs",
                "Cette étude a comparé plusieurs stratégies de filtrage de Kalman pour combiner les mesures IMU, GNSS et altitude barométrique. "
                "Le meilleur résultat a utilisé un réglage adaptatif de covariance pendant la dégradation du GNSS.",
            ),
            (
                "Comportement en vol",
                "La solution de navigation est restée stable pendant de courtes coupures GNSS lorsque le biais IMU était recalibré avant le décollage. "
                "Les longues coupures nécessitent encore un appui par cartographie ou odométrie visuelle.",
            ),
            (
                "Tests de robustesse",
                "Les essais ont couvert les manœuvres rapides, les environnements dégradés et les variations de capteurs. "
                "La qualité de la fusion dépendait fortement de la synchronisation temporelle et de la cohérence des offsets.",
            ),
            (
                "Suivi recommandé",
                "Le prochain prototype doit intégrer l'odométrie par vision par ordinateur, la comparer à la localisation par lidar et conserver un historique clair des dérives constatées.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_verification_docx(path: Path, scale: int) -> None:
    """Create a software verification document covering avionics constraints."""

    document = Document()
    _add_docx_sections(
        document,
        "Vérification logicielle avionique",
        [
            (
                "Contexte de certification",
                "Le lot logiciel visait une base de vérification adaptée aux contraintes avioniques, avec traçabilité des exigences, couverture de tests et journal des anomalies.",
            ),
            (
                "Méthode",
                "La méthode combinait tests unitaires, tests d'intégration et tests système sur banc. "
                "Chaque correction devait être associée à une exigence et à un résultat d'exécution archivés.",
            ),
            (
                "Résultats",
                "Les résultats montraient que la couverture fonctionnelle était suffisante, mais que certaines interfaces de bord restaient sensibles aux délais de communication et aux dépendances entre modules.",
            ),
            (
                "Gaps",
                "Les principaux écarts concernaient la robustesse face aux entrées invalides, la documentation incomplète de certains cas limites et le manque de scénarios de non-régression automatisés.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_thermal_docx(path: Path, scale: int) -> None:
    """Create a thermal management and power budget document."""

    document = Document()
    _add_docx_sections(
        document,
        "Gestion thermique et budget de puissance",
        [
            (
                "Objet",
                "Le document analyse la gestion thermique d'un banc électronique embarqué et le compromis entre performances de calcul et dissipation de chaleur.",
            ),
            (
                "Mesures",
                "Les mesures ont montré que les pics de puissance apparaissaient pendant les séquences de calcul intensif et les périodes d'acquisition capteur. "
                "Le refroidissement passif restait acceptable tant que la charge restait sous les seuils définis.",
            ),
            (
                "Risques",
                "Les risques principaux sont la surchauffe locale, le vieillissement accéléré des composants et la baisse de stabilité quand la température dépasse le cadre nominal.",
            ),
            (
                "Actions",
                "Les actions recommandées incluent l'amélioration de la conduction thermique, la surveillance continue des températures et la révision des scénarios de charge.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_cyber_docx(path: Path, scale: int) -> None:
    """Create a cybersecurity and data-link hardening document."""

    document = Document()
    _add_docx_sections(
        document,
        "Cybersécurité des liaisons de données",
        [
            (
                "Menaces",
                "Le document recense les menaces sur une liaison de données embarquée, notamment l'altération de paquets, la compromission de paramètres et la perte d'intégrité des journaux.",
            ),
            (
                "Mesures de protection",
                "Les protections combinent authentification, contrôle d'intégrité, segmentation réseau et journalisation sécurisée. "
                "Le modèle de menace distingue les accès autorisés, les tentatives d'injection et les attaques de rebond.",
            ),
            (
                "Validation",
                "La validation doit démontrer que les mécanismes de sécurité n'affectent pas la latence critique ni la disponibilité des services temps réel.",
            ),
            (
                "Suivi",
                "Le suivi recommandé impose des revues régulières, une mise à jour des clés et des tests de résistance sur les scénarios de perte de liaison.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_optronics_docx(path: Path, scale: int) -> None:
    """Create an optronics and image-processing document."""

    document = Document()
    _add_docx_sections(
        document,
        "Chaîne optronique et traitement d'image",
        [
            (
                "Contexte",
                "La chaîne optronique combine acquisition vidéo, calibration capteur et traitement d'image embarqué pour détecter des objets d'intérêt dans des scènes complexes.",
            ),
            (
                "Algorithmes",
                "Les algorithmes évalués incluent le recalage d'image, la détection de contours, la stabilisation et la décomposition du flux vidéo en plans exploitables.",
            ),
            (
                "Risques",
                "Les risques concernent la variabilité des scènes, la latence de calcul et la sensibilité aux vibrations et au bruit de capteur.",
            ),
            (
                "Suite",
                "La suite doit comparer l'exécution CPU à une accélération matérielle et documenter les écarts de performance en conditions nominales et dégradées.",
            ),
        ],
        scale,
    )
    document.save(path)


def create_uav_pptx(path: Path, scale: int) -> None:
    """Create a UAV navigation presentation with several dense slides."""

    presentation = Presentation()
    _add_pptx_slides(
        presentation,
        "Feuille de route navigation UAV",
        [
            (
                "Pile de navigation",
                "La pile de navigation UAV combine GNSS, IMU, vision par ordinateur et fusion de capteurs. "
                "La feuille de route privilégie un fonctionnement robuste dans des environnements GNSS dégradés.",
            ),
            (
                "Risques ouverts",
                "Les risques ouverts incluent la dérive de l'odométrie visuelle, les limites de calcul embarqué et l'effort de validation pour les cas de sécurité avionique.",
            ),
            (
                "Priorités de développement",
                "Les priorités incluent la validation en environnement dégradé, l'amélioration de la fusion multi-capteurs et la réduction de la dépendance aux signaux GNSS stables.",
            ),
            (
                "Critères d'acceptation",
                "Les critères d'acceptation demandent une répétabilité suffisante, une latence maîtrisée et des preuves documentaires associant chaque exigence aux résultats d'essai.",
            ),
        ],
        scale,
    )
    presentation.save(path)


def create_avionics_pptx(path: Path, scale: int) -> None:
    """Create an avionics systems engineering presentation."""

    presentation = Presentation()
    _add_pptx_slides(
        presentation,
        "Vue système avionique",
        [
            (
                "Architecture",
                "L'architecture avionique met en relation capteurs, calculateur embarqué, liaisons de données et fonctions de supervision. "
                "Les interfaces doivent rester déterministes malgré les variations de charge.",
            ),
            (
                "Validation",
                "La validation couvre l'intégration matérielle, les essais environnementaux et les scénarios de panne ou de dégradation partielle.",
            ),
            (
                "Gouvernance",
                "La gouvernance impose une traçabilité forte des exigences, des preuves et des écarts jusqu'à la clôture.",
            ),
            (
                "Décisions",
                "Les décisions d'architecture doivent documenter les compromis entre performance, sécurité, consommation et maintenabilité.",
            ),
        ],
        scale,
    )
    presentation.save(path)


def create_ai_docx(path: Path, scale: int) -> None:
    """Create a document about edge AI on embedded hardware."""

    document = Document()
    _add_docx_sections(
        document,
        "Traitement IA embarqué en périphérie",
        [
            (
                "Objectif",
                "L'étude évalue l'exécution de modèles IA à la périphérie, sur des calculateurs embarqués à ressources limitées et sans GPU dédié.",
            ),
            (
                "Contraintes",
                "Les contraintes principales sont la mémoire disponible, le temps de réponse, la consommation et la robustesse des entrées capteurs.",
            ),
            (
                "Résultats",
                "Les résultats montrent qu'une quantification prudente et une sélection stricte des couches permettent de garder une latence acceptable.",
            ),
            (
                "Risques",
                "Les risques concernent le surajustement sur des jeux de données limités, la dérive en exploitation et la difficulté d'auditer les décisions.",
            ),
        ],
        scale,
    )
    document.save(path)


def _write_corpus(scale: int) -> None:
    """Generate the full sample corpus with the chosen scale."""

    SAMPLE_DIR.mkdir(parents=True, exist_ok=True)
    _clear_existing_samples()
    create_radar_docx(SAMPLE_DIR / "etude_radar_fpga.docx", scale)
    create_radar_validation_docx(SAMPLE_DIR / "campagne_validation_radar.docx", scale)
    create_sensor_docx(SAMPLE_DIR / "fusion_capteurs_gnss.docx", scale)
    create_verification_docx(SAMPLE_DIR / "verification_logicielle_avionique.docx", scale)
    create_thermal_docx(SAMPLE_DIR / "gestion_thermique_puissance.docx", scale)
    create_cyber_docx(SAMPLE_DIR / "cybersecurite_liaison_donnees.docx", scale)
    create_optronics_docx(SAMPLE_DIR / "chaine_optronique_traitement_image.docx", scale)
    create_ai_docx(SAMPLE_DIR / "traitement_ia_embarque.docx", scale)
    create_uav_pptx(SAMPLE_DIR / "feuille_route_navigation_uav.pptx", scale)
    create_avionics_pptx(SAMPLE_DIR / "vue_systeme_avionique.pptx", scale)


def main() -> None:
    """Generate a French demo corpus for local testing."""

    parser = argparse.ArgumentParser(description="Generate French sample documents for local testing.")
    parser.add_argument("--scale", type=int, default=1, help="Repeat technical content to make the corpus larger.")
    args = parser.parse_args()
    scale = max(1, args.scale)
    _write_corpus(scale)
    print(f"Documents de démonstration écrits dans {SAMPLE_DIR} (scale={scale})")


if __name__ == "__main__":
    main()
